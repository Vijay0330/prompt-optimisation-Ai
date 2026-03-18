import os
import re
import json
import base64
import asyncio
import httpx
from io import BytesIO
from typing import Dict, Any, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Load prompt template ──────────────────────────────────────────────────────
_PROMPT_PATH = os.path.join(os.path.dirname(__file__), "..", "prompts", "ppt_prompt.txt")
with open(_PROMPT_PATH, "r") as f:
    _PPT_TEMPLATE = f.read()

MAX_RETRIES  = 3
RETRY_DELAYS = [5, 15, 30]

# ── Theme definitions ─────────────────────────────────────────────────────────
THEMES = {
    "ocean": {
        "bg":          RGBColor(10,  14,  40),   # deep navy
        "title_text":  RGBColor(255, 255, 255),
        "accent":      RGBColor(56,  189, 248),   # cyan
        "bullet_text": RGBColor(226, 232, 240),
        "slide_bg":    RGBColor(15,  23,  60),
        "bar":         RGBColor(56,  189, 248),
    },
    "midnight": {
        "bg":          RGBColor(8,   8,   20),
        "title_text":  RGBColor(255, 255, 255),
        "accent":      RGBColor(167, 139, 250),   # violet
        "bullet_text": RGBColor(216, 210, 254),
        "slide_bg":    RGBColor(15,  12,  35),
        "bar":         RGBColor(139, 92,  246),
    },
    "forest": {
        "bg":          RGBColor(5,   46,  22),
        "title_text":  RGBColor(255, 255, 255),
        "accent":      RGBColor(134, 239, 172),   # lime
        "bullet_text": RGBColor(220, 252, 231),
        "slide_bg":    RGBColor(6,   60,  28),
        "bar":         RGBColor(34,  197, 94),
    },
    "sunset": {
        "bg":          RGBColor(30,  10,  60),
        "title_text":  RGBColor(255, 255, 255),
        "accent":      RGBColor(251, 146, 60),    # orange
        "bullet_text": RGBColor(255, 237, 213),
        "slide_bg":    RGBColor(45,  15,  75),
        "bar":         RGBColor(249, 115, 22),
    },
    "corporate": {
        "bg":          RGBColor(248, 250, 252),
        "title_text":  RGBColor(15,  23,  42),
        "accent":      RGBColor(79,  70,  229),   # indigo
        "bullet_text": RGBColor(51,  65,  85),
        "slide_bg":    RGBColor(241, 245, 249),
        "bar":         RGBColor(99,  102, 241),
    },
}

FONT_SIZES = {
    "small":  {"title": 32, "subtitle": 18, "slide_title": 24, "bullet": 14},
    "medium": {"title": 38, "subtitle": 22, "slide_title": 28, "bullet": 16},
    "large":  {"title": 44, "subtitle": 26, "slide_title": 32, "bullet": 18},
}


def _build_url(model: str) -> str:
    api_key = os.getenv("GEMINI_API_KEY", "")
    if not api_key:
        raise ValueError("GEMINI_API_KEY is not set.")
    return (
        f"https://generativelanguage.googleapis.com/v1beta/models/"
        f"{model}:generateContent?key={api_key}"
    )


async def _call_gemini(prompt: str, model: str) -> Dict[str, Any]:
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": 0.5, "maxOutputTokens": 3000},
    }
    url = _build_url(model)
    last_error = None

    for attempt in range(MAX_RETRIES):
        try:
            async with httpx.AsyncClient(timeout=45.0) as client:
                response = await client.post(url, json=payload)
                if response.status_code == 429:
                    await asyncio.sleep(RETRY_DELAYS[attempt])
                    last_error = "Rate limited"
                    continue
                response.raise_for_status()
                data = response.json()

            raw = (
                data.get("candidates", [{}])[0]
                .get("content", {})
                .get("parts", [{}])[0]
                .get("text", "")
            )
            cleaned = re.sub(r"```(?:json)?", "", raw).strip().strip("`")
            return json.loads(cleaned)

        except httpx.HTTPStatusError as e:
            if e.response.status_code == 429:
                await asyncio.sleep(RETRY_DELAYS[min(attempt, 2)])
                last_error = str(e)
                continue
            raise

    raise Exception(f"Gemini rate limit after {MAX_RETRIES} retries. {last_error}")


def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int, bold: bool, color: RGBColor,
                  align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def _build_title_slide(prs, slide_data: dict, theme: dict, sizes: dict,
                       custom_header: Optional[str]):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, theme["bg"])

    W = prs.slide_width
    H = prs.slide_height

    # Accent bar left edge
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.12), H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["bar"]
    bar.line.fill.background()

    # Bottom accent strip
    strip = slide.shapes.add_shape(
        1, Inches(0), H - Inches(0.08), W, Inches(0.08),
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = theme["accent"]
    strip.line.fill.background()

    title_text = custom_header or slide_data.get("title", "Presentation")
    _add_text_box(
        slide, title_text,
        Inches(0.6), Inches(2.2), W - Inches(1.2), Inches(1.8),
        sizes["title"], True, theme["title_text"],
    )

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_text_box(
            slide, subtitle,
            Inches(0.6), Inches(3.8), W - Inches(1.2), Inches(1.0),
            sizes["subtitle"], False, theme["accent"],
        )

    _add_text_box(
        slide, "Prompt Intelligence Assistant",
        Inches(0.6), H - Inches(0.7), W - Inches(1.2), Inches(0.4),
        9, False, RGBColor(120, 120, 160),
    )


def _build_content_slide(prs, slide_data: dict, theme: dict, sizes: dict):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, theme["slide_bg"])

    W = prs.slide_width
    H = prs.slide_height

    # Top accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), W, Inches(0.07),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["bar"]
    bar.line.fill.background()

    # Slide title
    _add_text_box(
        slide, slide_data.get("title", ""),
        Inches(0.5), Inches(0.2), W - Inches(1.0), Inches(0.9),
        sizes["slide_title"], True, theme["accent"],
    )

    # Thin separator
    sep = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.05), W - Inches(1.0), Inches(0.03),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = theme["bar"]
    sep.line.fill.background()

    # Bullet points
    bullets = slide_data.get("bullets", [])
    y_start = Inches(1.2)
    row_h   = Inches(0.58)

    for i, bullet in enumerate(bullets[:5]):
        y = y_start + i * row_h

        # Bullet dot
        dot = slide.shapes.add_shape(
            1, Inches(0.5), y + Inches(0.15), Inches(0.12), Inches(0.12),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = theme["accent"]
        dot.line.fill.background()

        _add_text_box(
            slide, bullet,
            Inches(0.75), y, W - Inches(1.25), row_h,
            sizes["bullet"], False, theme["bullet_text"],
        )

    # Slide number (bottom right)
    _add_text_box(
        slide,
        str(slide_data.get("slide_number", "")),
        W - Inches(0.8), H - Inches(0.4), Inches(0.5), Inches(0.3),
        9, False, RGBColor(100, 100, 140),
        align=PP_ALIGN.RIGHT,
    )


def _build_pptx(doc: dict, theme_name: str, font_size: str,
                custom_header: Optional[str]) -> bytes:
    theme  = THEMES.get(theme_name, THEMES["ocean"])
    sizes  = FONT_SIZES.get(font_size, FONT_SIZES["medium"])

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_data in doc.get("slides", []):
        if slide_data.get("type") == "title":
            _build_title_slide(prs, slide_data, theme, sizes, custom_header)
        else:
            _build_content_slide(prs, slide_data, theme, sizes)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def generate_ppt(
    topic:         str,
    model:         str = "gemini-2.0-flash",
    theme:         str = "ocean",
    font_size:     str = "medium",
    custom_header: Optional[str] = None,
) -> Tuple[bytes, Dict]:
    """
    Main entry — generate a styled PPTX from a topic.
    Returns (pptx_bytes, metadata_dict).
    """
    prompt = _PPT_TEMPLATE.format(topic=topic)
    doc    = await _call_gemini(prompt, model)

    pptx_bytes = _build_pptx(doc, theme, font_size, custom_header)
    slide_count = len(doc.get("slides", []))

    metadata = {
        "title":       doc.get("title", "Presentation"),
        "slide_count": slide_count,
        "theme":       theme,
    }
    return pptx_bytes, metadata
